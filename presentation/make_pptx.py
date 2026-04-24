"""
Recreate the Beamer presentation as a .pptx (UNSW colour scheme).
Run:  python presentation/make_pptx.py
Output: presentation/presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── UNSW colours ─────────────────────────────────────────────────────────────
NAV   = RGBColor(0x00, 0x0F, 0x2B)   # #000F2B  dark navy
YEL   = RGBColor(0xFF, 0xD2, 0x00)   # #FFD200  yellow
RED   = RGBColor(0xE6, 0x3E, 0x30)   # red
GRN   = RGBColor(0x00, 0x84, 0x3D)   # green
LGREY = RGBColor(0xF0, 0xF0, 0xF0)   # light grey
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)


# ── low-level helpers ─────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank(prs):
    return prs.slide_layouts[6]

def rect(slide, l, t, w, h, rgb, line=False):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb
    if line:
        sh.line.color.rgb = rgb
        sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    return sh

def txb(slide, l, t, w, h, text, sz=14, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.size   = Pt(sz)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def para(tf, text, sz=13, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(sz)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p

# ── slide chrome: yellow title bar + progress dot + footer ────────────────────
BAR_H  = Inches(0.85)
FOOT_Y = Inches(7.15)
FOOT_H = Inches(0.35)

def title_bar(slide, title, section=""):
    rect(slide, 0, 0, SLIDE_W, BAR_H, YEL)
    # thin navy stripe at bottom of bar
    rect(slide, 0, BAR_H - Inches(0.04), SLIDE_W, Inches(0.04), NAV)
    txb(slide, Inches(0.3), Inches(0.10), Inches(12.5), Inches(0.65),
        title, sz=24, bold=True, color=NAV)

def footer(slide, section=""):
    rect(slide, 0, FOOT_Y, SLIDE_W, FOOT_H, NAV)
    foot = "Multi-Task Financial NLP  ·  COMP6713 2026 T1  ·  UNSW"
    if section:
        foot = section + "  ·  " + foot
    txb(slide, Inches(0.3), FOOT_Y + Inches(0.05), Inches(12.5), Inches(0.25),
        foot, sz=9, color=RGBColor(0xAA, 0xCC, 0xFF))

BODY_Y = BAR_H + Inches(0.10)
BODY_H = FOOT_Y - BAR_H - Inches(0.15)

# ── section divider slide ─────────────────────────────────────────────────────
def section_slide(prs, name):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, NAV)
    rect(sl, 0, Inches(3.4), SLIDE_W, Inches(0.10), YEL)
    txb(sl, Inches(1.0), Inches(2.8), Inches(11.0), Inches(1.2),
        name, sz=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    footer(sl, name)

# =============================================================================
# SLIDE 1 – Title
# =============================================================================
def slide_title(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, NAV)
    rect(sl, 0, Inches(5.05), SLIDE_W, Inches(0.10), YEL)

    txb(sl, Inches(0.8), Inches(0.7), Inches(11.5), Inches(1.8),
        "Multi-Task Financial NLP:\nStance & Sentiment Classification of\nFOMC Communications and News Headlines",
        sz=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, Inches(0.8), Inches(2.65), Inches(11.5), Inches(0.5),
        "COMP6713 — Natural Language Processing — 2026 T1",
        sz=16, color=YEL, align=PP_ALIGN.CENTER)
    txb(sl, Inches(0.8), Inches(3.35), Inches(11.5), Inches(0.45),
        "Louis Nguyen  ·  Quoc Dat Bui  ·  Nam Khanh Tran  ·  Quang Minh Phan",
        sz=15, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, Inches(0.8), Inches(3.85), Inches(11.5), Inches(0.4),
        "UNSW School of Computer Science and Engineering",
        sz=13, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 2 – Outline
# =============================================================================
def slide_outline(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Outline")
    footer(sl)

    sections = [
        "1.  Problem — Two financial NLP tasks and why they matter",
        "2.  Datasets — FOMC, Financial PhraseBank, Loughran–McDonald Lexicon",
        "3.  Methodology — 6-stage pipeline from lexicon rules to multi-task transformers",
        "4.  Results — 12 systems, two tasks, side-by-side comparison",
        "5.  Analysis — Key findings, domain-pretraining gap, performance progression",
        "6.  Deployment — CLI, Gradio demo, Hugging Face Hub release",
        "7.  Summary",
    ]
    box = slide.shapes.add_textbox if False else None
    tb  = sl.shapes.add_textbox(Inches(0.8), BODY_Y + Inches(0.2),
                                 Inches(11.5), BODY_H - Inches(0.2))
    tf  = tb.text_frame
    tf.word_wrap = True
    for i, s in enumerate(sections):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = s
        r.font.size  = Pt(16)
        r.font.color.rgb = NAV

# =============================================================================
# SLIDE 3 – Two Tasks
# =============================================================================
def slide_two_tasks(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Two Financial NLP Classification Tasks", "Problem")
    footer(sl, "Problem")

    txb(sl, Inches(0.4), BODY_Y + Inches(0.05), Inches(12.5), Inches(0.35),
        "We tackle two complementary 3-class classification problems over short financial text.",
        sz=13, color=NAV)

    # left card — Stance
    cy = BODY_Y + Inches(0.5)
    cw = Inches(5.9)
    ch = Inches(4.8)
    rect(sl, Inches(0.4), cy, cw, ch, LGREY, line=True)
    rect(sl, Inches(0.4), cy, cw, Inches(0.38), NAV)
    txb(sl, Inches(0.5), cy + Inches(0.04), cw - Inches(0.1), Inches(0.30),
        "Stance — FOMC Communications", sz=13, bold=True, color=WHITE)
    tb = sl.shapes.add_textbox(Inches(0.5), cy + Inches(0.45),
                                cw - Inches(0.15), ch - Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Label the monetary-policy stance of a sentence from U.S. Federal Reserve meetings (Shah et al., 2023).", sz=12)
    para(tf, "  • dovish — easing bias", sz=12, color=GRN)
    para(tf, "  • hawkish — tightening bias", sz=12, color=RED)
    para(tf, "  • neutral — no clear bias", sz=12, color=NAV)
    para(tf, "", sz=8)
    para(tf, '"…the Committee decided to lower the target range for the federal funds rate."  →  dovish',
         sz=11, italic=True, color=RGBColor(0x33, 0x33, 0x33))

    # right card — Sentiment
    rect(sl, Inches(7.0), cy, cw, ch, LGREY, line=True)
    rect(sl, Inches(7.0), cy, cw, Inches(0.38), NAV)
    txb(sl, Inches(7.1), cy + Inches(0.04), cw - Inches(0.1), Inches(0.30),
        "Sentiment — Financial PhraseBank", sz=13, bold=True, color=WHITE)
    tb2 = sl.shapes.add_textbox(Inches(7.1), cy + Inches(0.45),
                                  cw - Inches(0.15), ch - Inches(0.55))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    para(tf2, "Label the investor-facing tone of a news headline (Malo et al., 2014).", sz=12)
    para(tf2, "  • negative", sz=12, color=RED)
    para(tf2, "  • neutral", sz=12, color=NAV)
    para(tf2, "  • positive", sz=12, color=GRN)
    para(tf2, "", sz=8)
    para(tf2, '"Operating profit totalled EUR 21.1 mn, up from EUR 18.6 mn a year earlier."  →  positive',
         sz=11, italic=True, color=RGBColor(0x33, 0x33, 0x33))

# =============================================================================
# SLIDE 4 – Why It Matters
# =============================================================================
def slide_why_matters(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Why It Matters", "Problem")
    footer(sl, "Problem")

    items = [
        ("Central-bank stance moves markets.",
         "A single sentence shift between hawkish and dovish language in an FOMC statement can reprice rates, equities, and FX within seconds."),
        ("News sentiment drives short-horizon prices.",
         "Headline tone correlates with intraday returns and is a staple input to systematic trading and risk models."),
        ("Both tasks share domain vocabulary but differ in pragmatics:",
         "Stance is about policy direction; sentiment is about corporate outlook. This motivates joint modelling."),
        ("Labelled data is scarce and expensive.",
         "FOMC: 2,480 sentences; FPB allagree: 2,264 sentences. Techniques that squeeze more signal from small corpora are essential."),
        ("Practical goal:",
         "A single deployable encoder that does both jobs competitively."),
    ]
    y = BODY_Y + Inches(0.15)
    for title, body in items:
        rect(sl, Inches(0.3), y + Inches(0.09), Inches(0.08), Inches(0.28), YEL)
        tb = sl.shapes.add_textbox(Inches(0.55), y, Inches(12.0), Inches(0.5))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = title + " "; r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = NAV
        r2 = p.add_run(); r2.text = body;          r2.font.size = Pt(13); r2.font.color.rgb = BLACK
        y += Inches(0.56)

# =============================================================================
# SLIDE 5 – Dataset Overview
# =============================================================================
def slide_datasets(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Dataset Overview", "Datasets")
    footer(sl, "Datasets")

    # table header
    hy = BODY_Y + Inches(0.15)
    cols = [Inches(3.4), Inches(1.3), Inches(2.5), Inches(3.0), Inches(2.7)]
    headers = ["Name", "Sentences", "Classes", "Label distribution (test)", "Notes"]
    x = Inches(0.3)
    for hdr, w in zip(headers, cols):
        rect(sl, x, hy, w - Inches(0.04), Inches(0.38), NAV)
        txb(sl, x + Inches(0.05), hy + Inches(0.04), w - Inches(0.1), Inches(0.30),
            hdr, sz=11, bold=True, color=WHITE)
        x += w

    rows = [
        ("FOMC – Trillion Dollar Words\n(Shah et al., 2023)",
         "2,480", "3 (dovish / hawkish / neutral)",
         "dovish 130 (26.2%)\nhawkish 121 (24.4%)\nneutral 245 (49.4%)",
         "FOMC statements, minutes, press conferences. Neutral-heavy."),
        ("Financial PhraseBank allagree\n(Malo et al., 2014)",
         "2,264", "3 (neg / neutral / pos)",
         "negative 61 (13.5%)\nneutral 278 (61.4%)\npositive 114 (25.2%)",
         "News headlines with 100% annotator agreement. Neutral-dominant."),
        ("Loughran–McDonald Lexicon\n(Loughran & McDonald, 2011)",
         "—", "rule-based\n(pos / neg / uncertainty)",
         "pos ~110 terms\nneg ~230 terms\nuncertainty ~90 terms",
         "Financial word lists; also hawkish (~40) & dovish (~50) monetary-policy lists."),
    ]
    alt = [WHITE, LGREY, WHITE]
    row_h = Inches(1.18)
    for ri, (row, bg) in enumerate(zip(rows, alt)):
        ry = hy + Inches(0.38) + ri * row_h
        x  = Inches(0.3)
        for cell, w in zip(row, cols):
            rect(sl, x, ry, w - Inches(0.04), row_h, bg, line=True)
            txb(sl, x + Inches(0.05), ry + Inches(0.04),
                w - Inches(0.12), row_h - Inches(0.08),
                cell, sz=10, color=NAV)
            x += w

    y_note = hy + Inches(0.38) + len(rows) * row_h + Inches(0.1)
    txb(sl, Inches(0.3), y_note, Inches(12.5), Inches(0.55),
        "• Short texts (mostly < 40 tokens) from the financial domain.\n"
        "• FOMC and FPB are imbalanced with dominant neutral; LM Lexicon provides a domain-aware rule-based baseline.",
        sz=11, color=NAV)

# =============================================================================
# SLIDE 6 – Preprocessing & Splits
# =============================================================================
def slide_preprocessing(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Preprocessing & Splits", "Datasets")
    footer(sl, "Datasets")

    # left column
    tb = sl.shapes.add_textbox(Inches(0.4), BODY_Y + Inches(0.1),
                                Inches(6.8), BODY_H - Inches(0.1))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Splits (stratified, random_state=42):", sz=13, bold=True, color=NAV)
    para(tf, "  • 70% train / 10% validation / 20% test", sz=12)
    para(tf, "  • FOMC: 1736 / 248 / 496", sz=12)
    para(tf, "  • FPB allagree: 1584 / 227 / 453", sz=12)
    para(tf, "", sz=7)
    para(tf, "Text normalisation:", sz=13, bold=True, color=NAV)
    para(tf, "  • Unicode NFKC, strip control chars, collapse whitespace", sz=12)
    para(tf, "  • Lower-case only for TF-IDF baselines", sz=12)
    para(tf, "  • Keep case & punctuation for BERT-family (max_len=128)", sz=12)
    para(tf, "", sz=7)
    para(tf, "Class weights for FOMC (weighted cross-entropy):", sz=13, bold=True, color=NAV)
    para(tf, "  w_c = N / (K · n_c)     →     w = (1.272, 1.365, 0.675)", sz=12)
    para(tf, "  for (dovish, hawkish, neutral) with K=3, N=1736", sz=11, color=RGBColor(0x44, 0x44, 0x44))

    # right column — note about chart
    rect(sl, Inches(7.6), BODY_Y + Inches(0.1), Inches(5.3), BODY_H - Inches(0.1), LGREY, line=True)
    txb(sl, Inches(7.7), BODY_Y + Inches(0.2), Inches(5.1), Inches(0.4),
        "Class Distribution", sz=13, bold=True, color=NAV)
    txb(sl, Inches(7.7), BODY_Y + Inches(0.7), Inches(5.1), Inches(3.5),
        "[ analysis/class_distribution.png ]\n\n"
        "Neutral dominates both tasks.\n"
        "Negative (FPB) and hawkish (FOMC) are the scarcest labels.",
        sz=12, color=RGBColor(0x44, 0x44, 0x44))

# =============================================================================
# SLIDE 7 – Pipeline Overview
# =============================================================================
def slide_pipeline(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Pipeline Overview", "Methodology")
    footer(sl, "Methodology")

    stages = [
        ("1. LM Lexicon",      "rule-based",             YEL),
        ("2. TF-IDF",          "+LR / +SVM",            YEL),
        ("3. Zero-/Few-shot",  "probes",                 RGBColor(0xFF, 0xE5, 0xE0)),
        ("4. FinBERT",         "fine-tune",              RGBColor(0xFF, 0xE5, 0xE0)),
        ("5. BERT LLRD",       "+ gradual UF",           LGREY),
    ]

    bx = Inches(0.4)
    by = BODY_Y + Inches(0.3)
    bw = Inches(2.2)
    bh = Inches(1.2)
    gap = Inches(0.25)

    for i, (title, sub, bg) in enumerate(stages):
        x = bx + i * (bw + gap)
        rect(sl, x, by, bw, bh, bg, line=True)
        txb(sl, x + Inches(0.1), by + Inches(0.15), bw - Inches(0.2), Inches(0.5),
            title, sz=12, bold=True, color=NAV)
        txb(sl, x + Inches(0.1), by + Inches(0.65), bw - Inches(0.2), Inches(0.35),
            sub, sz=11, color=RGBColor(0x44, 0x44, 0x44))
        if i < len(stages) - 1:
            ax = x + bw + Inches(0.04)
            ay = by + bh / 2 - Inches(0.04)
            txb(sl, ax, ay, gap - Inches(0.04), Inches(0.25), "→", sz=18, bold=True, color=NAV, align=PP_ALIGN.CENTER)

    # multi-task box below
    mtx = Inches(4.7)
    mty = by + bh + Inches(0.6)
    mtw = Inches(4.0)
    mth = Inches(1.0)
    rect(sl, mtx, mty, mtw, mth, RGBColor(0xCC, 0xEE, 0xD8), line=True)
    txb(sl, mtx + Inches(0.1), mty + Inches(0.1), mtw - Inches(0.2), Inches(0.5),
        "6. Multi-task FinBERT", sz=13, bold=True, color=NAV)
    txb(sl, mtx + Inches(0.1), mty + Inches(0.58), mtw - Inches(0.2), Inches(0.35),
        "shared encoder + 2 heads", sz=11, color=RGBColor(0x44, 0x44, 0x44))

    y_note = mty + mth + Inches(0.2)
    tb = sl.shapes.add_textbox(Inches(0.4), y_note, Inches(12.5), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Each stage is a controlled upgrade: add features → add domain prior → add fine-tuning → add training recipe → add multi-task structure.", sz=12, color=NAV)
    para(tf, "Every system is evaluated on the same held-out test sets — gains are directly comparable.", sz=12, color=NAV)

# =============================================================================
# SLIDE 8 – Baselines
# =============================================================================
def slide_baselines(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Baselines: Lexicon and TF-IDF", "Methodology")
    footer(sl, "Methodology")

    items = [
        ("LM Lexicon (rule-based) [Loughran & McDonald 2011]",
         "Count positive/negative terms per sentence, predict by sign. Domain-specific but brittle — sign-counting cannot handle negation or context."),
        ("TF-IDF + Logistic Regression",
         "Bigrams (ngram=(1,2)), max_features=50k, sublinear_tf=True; LR with default ℓ₂ regularisation."),
        ("TF-IDF + Linear SVM",
         "Same vectoriser; LinearSVC, C=1.0, class_weight='balanced' to offset the neutral skew."),
        ("TF-IDF (trigrams) + LR",
         "Richer context: ngram=(1,3), max_features=80k, min_df=2. Adds monetary-policy phrases like 'target range federal funds rate'."),
        ("TF-IDF + LM Lexicon (hybrid) + LR",
         "Concatenate TF-IDF vectors with 8 normalised LM lexicon features (counts ÷ total words + net polarity). Tests whether hand-crafted signals complement surface n-grams."),
    ]
    y = BODY_Y + Inches(0.15)
    for title, body in items:
        rect(sl, Inches(0.3), y + Inches(0.08), Inches(0.08), Inches(0.28), YEL)
        tb = sl.shapes.add_textbox(Inches(0.55), y, Inches(12.3), Inches(0.55))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = title + " — "; r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = NAV
        r2 = p.add_run(); r2.text = body;             r2.font.size = Pt(12); r2.font.color.rgb = BLACK
        y += Inches(0.62)

    txb(sl, Inches(0.3), y + Inches(0.1), Inches(12.5), Inches(0.4),
        "These establish a domain-aware but shallow floor. Any neural model must clearly beat them.",
        sz=11, italic=True, color=RGBColor(0x44, 0x44, 0x44))

# =============================================================================
# SLIDE 9 – Pre-trained evaluation
# =============================================================================
def slide_pretrained(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Pre-trained Evaluation: Zero-shot and Few-shot", "Methodology")
    footer(sl, "Methodology")

    # left panel
    rect(sl, Inches(0.3), BODY_Y + Inches(0.1), Inches(6.0), BODY_H - Inches(0.1), LGREY, line=True)
    txb(sl, Inches(0.5), BODY_Y + Inches(0.22), Inches(5.7), Inches(0.35),
        "Zero-shot FinBERT (native 3-way head)", sz=13, bold=True, color=NAV)
    tb = sl.shapes.add_textbox(Inches(0.5), BODY_Y + Inches(0.65), Inches(5.8), Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "  • Use the released {positive, negative, neutral} head directly on FPB — no parameter updates.", sz=12)
    para(tf, "  • For FOMC: proxy map  positive → hawkish,  negative → dovish,  neutral → neutral.", sz=12)
    para(tf, "    Tests whether sentiment is a usable proxy for stance.", sz=12)
    para(tf, "  • No parameter updates.", sz=12)

    # right panel
    rect(sl, Inches(6.9), BODY_Y + Inches(0.1), Inches(6.0), BODY_H - Inches(0.1), LGREY, line=True)
    txb(sl, Inches(7.1), BODY_Y + Inches(0.22), Inches(5.7), Inches(0.35),
        "Few-shot linear probe (k=16 per class)", sz=13, bold=True, color=NAV)
    tb2 = sl.shapes.add_textbox(Inches(7.1), BODY_Y + Inches(0.65), Inches(5.7), Inches(4.5))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    para(tf2, "  • Freeze encoder, extract [CLS] embedding, fit logistic regression on 3 × 16 = 48 examples.", sz=12)
    para(tf2, "  • Compare three frozen backbones:", sz=12)
    para(tf2, "      – FinBERT  (ProsusAI/finbert)", sz=12, color=NAV)
    para(tf2, "      – BERT-base-uncased  (Devlin et al., 2019)", sz=12, color=NAV)
    para(tf2, "      – RoBERTa-base", sz=12, color=NAV)
    para(tf2, "  • Isolates domain pretraining value in a tiny-data setting.", sz=12)

# =============================================================================
# SLIDE 10 – Fine-tuning
# =============================================================================
def slide_finetuning(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Fine-tuning: FinBERT vs. BERT-base LLRD", "Methodology")
    footer(sl, "Methodology")

    # left block
    rect(sl, Inches(0.3), BODY_Y + Inches(0.1), Inches(6.0), BODY_H - Inches(0.2), NAV)
    txb(sl, Inches(0.5), BODY_Y + Inches(0.22), Inches(5.6), Inches(0.38),
        "FinBERT single-task fine-tune", sz=13, bold=True, color=YEL)
    tb = sl.shapes.add_textbox(Inches(0.5), BODY_Y + Inches(0.65), Inches(5.6), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    for item in [
        "Backbone: ProsusAI/finbert",
        "Optimiser: AdamW, lr 2e-5, weight decay 0.01",
        "Schedule: linear warmup (10%) + linear decay",
        "5 epochs, batch size 32, max_len=128",
        "Weighted cross-entropy for FOMC",
        "Strong domain prior, conventional recipe.",
    ]:
        para(tf, "  • " + item, sz=12, color=WHITE)

    # right block
    rect(sl, Inches(6.9), BODY_Y + Inches(0.1), Inches(6.0), BODY_H - Inches(0.2), NAV)
    txb(sl, Inches(7.1), BODY_Y + Inches(0.22), Inches(5.6), Inches(0.38),
        "BERT-base LLRD + Gradual Unfreeze", sz=13, bold=True, color=YEL)
    tb2 = sl.shapes.add_textbox(Inches(7.1), BODY_Y + Inches(0.65), Inches(5.6), Inches(5.0))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for item in [
        "Backbone: bert-base-uncased",
        "Layer-wise LR decay: head 2e-5, × 0.9 per depth  →  η_d = 2×10⁻⁵ · 0.9^d",
        "Label smoothing ε = 0.1",
        "10 epochs, batch size 32",
        "Gradual unfreeze: 1 layer group unlocked per epoch",
        "General-purpose backbone — no financial pretraining.",
    ]:
        para(tf2, "  • " + item, sz=12, color=WHITE)

    txb(sl, Inches(0.3), FOOT_Y - Inches(0.45), Inches(12.5), Inches(0.4),
        "Key comparison: does domain pretraining (FinBERT) beat a better fine-tuning recipe (BERT-base)?  Result: they tie.",
        sz=12, italic=True, color=NAV)

# =============================================================================
# SLIDE 11 – Multi-task Architecture
# =============================================================================
def slide_multitask(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Multi-task FinBERT Architecture", "Methodology")
    footer(sl, "Methodology")

    # architecture diagram (textbox art)
    diag_x = Inches(0.6)
    diag_y = BODY_Y + Inches(0.15)

    # Input boxes
    rect(sl, diag_x,              diag_y,               Inches(2.5), Inches(0.7), YEL, line=True)
    rect(sl, diag_x + Inches(3.2), diag_y,               Inches(2.5), Inches(0.7), YEL, line=True)
    txb(sl, diag_x + Inches(0.1), diag_y + Inches(0.15), Inches(2.3), Inches(0.4),
        "FOMC batch", sz=12, bold=True, color=NAV, align=PP_ALIGN.CENTER)
    txb(sl, diag_x + Inches(3.3), diag_y + Inches(0.15), Inches(2.3), Inches(0.4),
        "FPB batch",  sz=12, bold=True, color=NAV, align=PP_ALIGN.CENTER)

    # Shared encoder
    enc_y = diag_y + Inches(1.0)
    rect(sl, diag_x, enc_y, Inches(5.7), Inches(0.85), LGREY, line=True)
    txb(sl, diag_x + Inches(0.1), enc_y + Inches(0.08), Inches(5.5), Inches(0.5),
        "Shared FinBERT Encoder  →  768-dim [CLS]", sz=13, bold=True, color=NAV, align=PP_ALIGN.CENTER)

    # Dropout
    drop_y = enc_y + Inches(0.85)
    rect(sl, diag_x + Inches(1.0), drop_y, Inches(3.7), Inches(0.55), WHITE, line=True)
    txb(sl, diag_x + Inches(1.0), drop_y + Inches(0.08), Inches(3.7), Inches(0.35),
        "Dropout(0.1)", sz=12, color=NAV, align=PP_ALIGN.CENTER)

    # Heads
    head_y = drop_y + Inches(0.9)
    rect(sl, diag_x,               head_y, Inches(2.5), Inches(0.75), RGBColor(0xFF, 0xE5, 0xE0), line=True)
    rect(sl, diag_x + Inches(3.2), head_y, Inches(2.5), Inches(0.75), RGBColor(0xCC, 0xEE, 0xD8), line=True)
    txb(sl, diag_x + Inches(0.1), head_y + Inches(0.05), Inches(2.3), Inches(0.6),
        "Stance head\nLinear(768, 3)", sz=11, bold=True, color=NAV, align=PP_ALIGN.CENTER)
    txb(sl, diag_x + Inches(3.3), head_y + Inches(0.05), Inches(2.3), Inches(0.6),
        "Sentiment head\nLinear(768, 3)", sz=11, bold=True, color=NAV, align=PP_ALIGN.CENTER)

    # Arrows (text approximation)
    for ax in [diag_x + Inches(1.1), diag_x + Inches(4.3)]:
        txb(sl, ax, diag_y + Inches(0.7), Inches(0.5), Inches(0.35), "↓", sz=18, color=NAV, align=PP_ALIGN.CENTER)
    txb(sl, diag_x + Inches(1.1), enc_y + Inches(0.85), Inches(0.5), Inches(0.35), "↓", sz=18, color=NAV, align=PP_ALIGN.CENTER)
    txb(sl, diag_x + Inches(0.8), drop_y + Inches(0.55), Inches(0.5), Inches(0.4), "↙", sz=18, color=NAV)
    txb(sl, diag_x + Inches(3.8), drop_y + Inches(0.55), Inches(0.5), Inches(0.4), "↘", sz=18, color=NAV)

    # right column — training details
    tb = sl.shapes.add_textbox(Inches(7.2), BODY_Y + Inches(0.1), Inches(5.7), BODY_H - Inches(0.1))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Training details (Liu et al., 2019):", sz=13, bold=True, color=NAV)
    for item in [
        "Alternating batches: one FOMC batch → one FPB batch; each updates its own head + shared encoder.",
        "AdamW, lr 2e-5, 8 epochs, batch size 32.",
        "Per-task cross-entropy (weighted CE for stance).",
        "Best checkpoint by average validation macro-F1 across both tasks.",
        "Apple M3 Max (MPS), seed 42.",
    ]:
        para(tf, "  • " + item, sz=12)
    para(tf, "", sz=7)
    para(tf, "Shared encoder ⇒ stance examples regularise the sentiment head (and vice versa) — a free data-augmentation effect.",
         sz=12, italic=True, color=RGBColor(0x33, 0x33, 0x33))

# =============================================================================
# SLIDE 12 – Baseline Results
# =============================================================================
def _table(slide, left, top, col_widths, headers, rows, row_highlights=()):
    """Draw a simple table. row_highlights = set of 0-based row indices to shade."""
    ROW_H = Inches(0.38)
    # header
    x = left
    for hdr, w in zip(headers, col_widths):
        rect(slide, x, top, w - Inches(0.03), ROW_H, NAV)
        txb(slide, x + Inches(0.06), top + Inches(0.06), w - Inches(0.1), ROW_H - Inches(0.08),
            hdr, sz=11, bold=True, color=WHITE)
        x += w
    y = top + ROW_H
    for ri, row in enumerate(rows):
        if row[0].startswith("§"):
            # group label
            label = row[0][1:]
            total_w = sum(col_widths)
            rect(slide, left, y, total_w, ROW_H - Inches(0.08), RGBColor(0xDD, 0xE8, 0xF5))
            txb(slide, left + Inches(0.08), y + Inches(0.04), total_w - Inches(0.1), ROW_H,
                label, sz=11, bold=True, color=NAV, italic=True)
            y += ROW_H - Inches(0.06)
            continue
        bg = YEL if ri in row_highlights else (LGREY if ri % 2 == 0 else WHITE)
        x = left
        for cell, w in zip(row, col_widths):
            rect(slide, x, y, w - Inches(0.03), ROW_H, bg, line=True)
            txb(slide, x + Inches(0.06), y + Inches(0.04), w - Inches(0.1), ROW_H - Inches(0.08),
                cell, sz=10, color=NAV)
            x += w
        y += ROW_H
    return y

def slide_baseline_results(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Baseline & Lexicon Results", "Results")
    footer(sl, "Results")

    cols = [Inches(4.2), Inches(1.6), Inches(1.5), Inches(1.7), Inches(1.5)]
    hdrs = ["Model", "Sent. Acc", "Sent. F1", "Stance Acc", "Stance F1"]
    rows = [
        ("LM Lexicon (rule-based)", "0.6932", "0.5315", "0.4153", "0.3885"),
        ("TF-IDF + LR",             "0.8720", "0.8232", "0.6089", "0.5873"),
        ("TF-IDF + SVM (LinearSVC)","0.8940★","0.8534★","0.6331★","0.6061★"),
        ("TF-IDF (trigrams) + LR",  "0.8786", "0.8310", "0.6109", "0.5914"),
        ("TF-IDF + LM Lexicon",     "0.8543", "0.8050", "0.6109", "0.5863"),
    ]
    end_y = _table(sl, Inches(0.3), BODY_Y + Inches(0.15), cols, hdrs, rows, row_highlights={2})

    tb = sl.shapes.add_textbox(Inches(0.3), end_y + Inches(0.15), Inches(12.5), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    for note in [
        "TF-IDF + LinearSVM is the strongest non-neural baseline on both tasks.",
        "Lexicon alone is weak (FPB F1 0.5315, stance F1 0.3885) — sign-counting cannot handle negation or context.",
        "Trigrams give only marginal lift over bigrams (+0.5 pp) — vocabulary grows faster than signal.",
        "Adding LM-lexicon features to TF-IDF does not help — TF-IDF vocabulary already covers those words.",
    ]:
        para(tf, "  • " + note, sz=12, color=NAV)

# =============================================================================
# SLIDE 13 – Pre-trained Results
# =============================================================================
def slide_pretrained_results(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Pre-trained Models: Zero-shot and Few-shot", "Results")
    footer(sl, "Results")

    cols = [Inches(4.2), Inches(1.6), Inches(1.5), Inches(1.7), Inches(1.5)]
    hdrs = ["Model", "Sent. Acc", "Sent. F1", "Stance Acc", "Stance F1"]
    rows = [
        ("FinBERT zero-shot (native head)", "0.9735", "0.9650", "0.4980", "0.4874"),
        ("FinBERT few-shot (k=16 probe)",   "0.9779★","0.9670★","0.4859", "0.4534"),
        ("BERT-base few-shot (k=16)",       "0.7417", "0.6500", "0.3851", "0.3744"),
        ("RoBERTa-base few-shot (k=16)",    "0.7682", "0.6722", "0.3730", "0.3600"),
    ]
    end_y = _table(sl, Inches(0.3), BODY_Y + Inches(0.15), cols, hdrs, rows, row_highlights={1})

    tb = sl.shapes.add_textbox(Inches(0.3), end_y + Inches(0.15), Inches(12.5), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    for note in [
        "FinBERT zero-shot is already near-optimal on sentiment (F1 0.9650) — its pretraining matches FPB almost exactly.",
        "FinBERT dominates both probes: +30 pp sentiment F1 over BERT/RoBERTa at k=16 — domain pretraining, not size, is what matters in low-data regimes.",
        "Stance is much harder zero-shot: the positive→hawkish / negative→dovish proxy map only reaches F1 0.4874 — stance ≠ sentiment.",
        "The probe tops sentiment but loses a little stance F1 vs. native head — 48 examples is too few to learn stance.",
    ]:
        para(tf, "  • " + note, sz=12, color=NAV)

# =============================================================================
# SLIDE 14 – Fine-tuned Results
# =============================================================================
def slide_finetuned_results(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Fine-tuned Single-Task Models", "Results")
    footer(sl, "Results")

    cols = [Inches(4.8), Inches(1.6), Inches(1.5), Inches(1.7), Inches(1.5)]
    hdrs = ["Model", "Sent. Acc", "Sent. F1", "Stance Acc", "Stance F1"]
    rows = [
        ("FinBERT (fine-tuned)",          "0.9669", "0.9459", "0.6371", "0.6194"),
        ("BERT-base LLRD + Gradual UF ★", "0.9757", "0.9670★","0.6512★","0.6383★"),
    ]
    end_y = _table(sl, Inches(0.3), BODY_Y + Inches(0.3), cols, hdrs, rows, row_highlights={1})

    tb = sl.shapes.add_textbox(Inches(0.3), end_y + Inches(0.3), Inches(12.5), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    for note in [
        "BERT-base with LLRD + gradual unfreeze matches FinBERT on sentiment (macro-F1 tied at 0.9670) "
        "and exceeds it on stance by +1.89 pp F1.",
        "Takeaway: a disciplined fine-tuning recipe on a general-purpose backbone can close the domain-pretraining gap.",
        "First evidence in our study that 'domain BERT' is not the only path to strong finance-domain results.",
    ]:
        para(tf, "  • " + note, sz=13, color=NAV)

# =============================================================================
# SLIDE 15 – Multi-task Headline
# =============================================================================
def slide_multitask_result(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Multi-task FinBERT — Headline Result", "Results")
    footer(sl, "Results")

    # big numbers
    hy = BODY_Y + Inches(0.3)
    for i, (label, val1, val2) in enumerate([
        ("Sentiment", "Accuracy  0.9779", "Macro-F1  0.9666"),
        ("Stance",    "Accuracy  0.6492", "Macro-F1  0.6384"),
    ]):
        bx = Inches(0.5) + i * Inches(5.8)
        rect(sl, bx, hy, Inches(5.4), Inches(2.2), LGREY, line=True)
        txb(sl, bx + Inches(0.15), hy + Inches(0.1), Inches(5.1), Inches(0.42),
            label, sz=18, bold=True, color=NAV, align=PP_ALIGN.CENTER)
        txb(sl, bx + Inches(0.15), hy + Inches(0.6), Inches(5.1), Inches(0.55),
            val1, sz=22, bold=True, color=GRN, align=PP_ALIGN.CENTER)
        txb(sl, bx + Inches(0.15), hy + Inches(1.25), Inches(5.1), Inches(0.55),
            val2, sz=22, bold=True, color=GRN, align=PP_ALIGN.CENTER)

    # deltas
    dy = hy + Inches(2.4)
    rect(sl, Inches(0.5), dy, Inches(12.0), Inches(0.75), RGBColor(0xFF, 0xF7, 0xCC))
    txb(sl, Inches(0.7), dy + Inches(0.12), Inches(11.5), Inches(0.5),
        "Δ vs. single-task FinBERT:   Sentiment +2.07 pp F1   |   Stance +1.90 pp F1",
        sz=16, bold=True, color=NAV, align=PP_ALIGN.CENTER)

    # bullets below
    tb = sl.shapes.add_textbox(Inches(0.5), dy + Inches(0.9), Inches(12.0), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    for note in [
        "One shared encoder + two heads + alternating batches → improves both tasks over single-task FinBERT.",
        "Sentiment F1 reaches 0.9666 (matching BERT LLRD's 0.9670); stance F1 reaches 0.6384 (within 0.0001 of BERT LLRD).",
        "Free lunch: no new data, only a second 768×3 head — training cost comparable to single-task.",
    ]:
        para(tf, "  • " + note, sz=13, color=NAV)

# =============================================================================
# SLIDE 16 – Full Results Summary
# =============================================================================
def slide_full_results(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Results Summary — All 12 Systems", "Results")
    footer(sl, "Results")

    cols = [Inches(4.0), Inches(1.6), Inches(1.4), Inches(1.6), Inches(1.4)]
    hdrs = ["Model", "Sent. Acc", "Sent. F1", "Stance Acc", "Stance F1"]
    rows_raw = [
        ("§Traditional baselines",),
        ("LM Lexicon (rule-based)",       "0.6932","0.5315","0.4153","0.3885"),
        ("TF-IDF + LR",                   "0.8720","0.8232","0.6089","0.5873"),
        ("TF-IDF + SVM (LinearSVC)",      "0.8940","0.8534","0.6331","0.6061"),
        ("TF-IDF (trigrams) + LR",        "0.8786","0.8310","0.6109","0.5914"),
        ("TF-IDF + LM Lexicon (hybrid)",  "0.8543","0.8050","0.6109","0.5863"),
        ("§Pre-trained (zero-/few-shot)",),
        ("FinBERT zero-shot",             "0.9735","0.9650","0.4980","0.4874"),
        ("FinBERT few-shot (k=16)",       "0.9779","0.9670","0.4859","0.4534"),
        ("BERT-base few-shot (k=16)",     "0.7417","0.6500","0.3851","0.3744"),
        ("RoBERTa-base few-shot (k=16)",  "0.7682","0.6722","0.3730","0.3600"),
        ("§Fine-tuned",),
        ("FinBERT (fine-tuned)",          "0.9669","0.9459","0.6371","0.6194"),
        ("BERT-base LLRD + Gradual UF ★","0.9757","0.9670★","0.6512★","0.6383★"),
        ("§Multi-task",),
        ("Multi-task FinBERT ★",         "0.9779★","0.9666★","0.6492","0.6384★"),
    ]
    best_rows = {12, 14}  # 0-based after group rows are skipped in highlight calc — just colour the table manually
    end_y = _table(sl, Inches(0.3), BODY_Y + Inches(0.05), cols, hdrs, rows_raw)

    txb(sl, Inches(0.3), end_y + Inches(0.1), Inches(12.5), Inches(0.5),
        "Statistical tie at the top: Multi-task FinBERT and BERT-base LLRD are indistinguishable within ±0.1 pp on both macro-F1s. "
        "Two very different routes reach the same ceiling.",
        sz=11, italic=True, color=NAV)

# =============================================================================
# SLIDE 17 – Error Analysis
# =============================================================================
def slide_errors(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Error Analysis — Multi-task FinBERT", "Results")
    footer(sl, "Results")

    # Stance table
    rect(sl, Inches(0.3), BODY_Y + Inches(0.1), Inches(5.5), Inches(0.38), NAV)
    txb(sl, Inches(0.4), BODY_Y + Inches(0.16), Inches(5.3), Inches(0.28),
        "Stance (FOMC test, N=496)", sz=12, bold=True, color=WHITE)
    for ri, (cls, f1, bg) in enumerate([
        ("dovish",  "0.6174", LGREY),
        ("hawkish", "0.5869", WHITE),
        ("neutral (best)", "0.7109", YEL),
    ]):
        ty = BODY_Y + Inches(0.48) + ri * Inches(0.38)
        rect(sl, Inches(0.3), ty, Inches(3.5), Inches(0.36), bg, line=True)
        rect(sl, Inches(3.8), ty, Inches(2.0), Inches(0.36), bg, line=True)
        txb(sl, Inches(0.4), ty + Inches(0.06), Inches(3.3), Inches(0.28), cls, sz=11, color=NAV)
        txb(sl, Inches(3.9), ty + Inches(0.06), Inches(1.8), Inches(0.28), f1, sz=11, bold=True, color=NAV)

    tb = sl.shapes.add_textbox(Inches(0.3), BODY_Y + Inches(1.7), Inches(5.5), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Most errors sit on the neutral vs. leaning boundary — hedged language is genuinely ambiguous.", sz=11)
    para(tf, "Hawkish is the weakest class (lowest count in train, most lexically diverse).", sz=11)

    # Sample errors stance
    txb(sl, Inches(0.3), BODY_Y + Inches(3.3), Inches(5.5), Inches(0.32),
        "Sample misclassifications:", sz=12, bold=True, color=NAV)
    errors_s = [
        ('"…the new statement conveys the Committee…"', "gold hawkish → predicted dovish"),
        ('"Looking ahead, reports from retailer contacts…"', "gold neutral → predicted dovish"),
        ('"Longer-term inflation expectations have moved…"', "gold hawkish → predicted neutral"),
    ]
    y = BODY_Y + Inches(3.65)
    for quote, label in errors_s:
        txb(sl, Inches(0.4), y, Inches(5.3), Inches(0.27), quote, sz=10, italic=True, color=BLACK)
        txb(sl, Inches(0.4), y + Inches(0.27), Inches(5.3), Inches(0.22), label, sz=10, bold=True, color=RED)
        y += Inches(0.56)

    # Sentiment table
    rect(sl, Inches(7.0), BODY_Y + Inches(0.1), Inches(5.9), Inches(0.38), NAV)
    txb(sl, Inches(7.1), BODY_Y + Inches(0.16), Inches(5.7), Inches(0.28),
        "Sentiment (FPB test, N=453)", sz=12, bold=True, color=WHITE)
    for ri, (cls, f1, bg) in enumerate([
        ("negative", "0.9500", LGREY),
        ("neutral (best)", "0.9928", YEL),
        ("positive", "0.9569", WHITE),
    ]):
        ty = BODY_Y + Inches(0.48) + ri * Inches(0.38)
        rect(sl, Inches(7.0), ty, Inches(3.5), Inches(0.36), bg, line=True)
        rect(sl, Inches(10.5), ty, Inches(2.4), Inches(0.36), bg, line=True)
        txb(sl, Inches(7.1), ty + Inches(0.06), Inches(3.3), Inches(0.28), cls, sz=11, color=NAV)
        txb(sl, Inches(10.6), ty + Inches(0.06), Inches(2.2), Inches(0.28), f1, sz=11, bold=True, color=NAV)

    tb2 = sl.shapes.add_textbox(Inches(7.0), BODY_Y + Inches(1.7), Inches(5.9), Inches(1.6))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    para(tf2, "Only ~10 errors across 453 test sentences — sentiment is near-saturated.", sz=11)
    para(tf2, "Remaining errors are subtle 'compared with' comparative constructions.", sz=11)

    txb(sl, Inches(7.0), BODY_Y + Inches(3.3), Inches(5.9), Inches(0.32),
        "Sample misclassifications:", sz=12, bold=True, color=NAV)
    errors_f = [
        ('"Operating loss was EUR 179mn, compared to a loss…"', "gold positive → predicted negative"),
        ('"In Jan–Jun 2010, diluted loss per share…"', "gold negative → predicted positive"),
        ('"Unit costs for flight operations fell by 6.4 percent…"', "gold positive → predicted negative"),
    ]
    y2 = BODY_Y + Inches(3.65)
    for quote, label in errors_f:
        txb(sl, Inches(7.1), y2, Inches(5.7), Inches(0.27), quote, sz=10, italic=True, color=BLACK)
        txb(sl, Inches(7.1), y2 + Inches(0.27), Inches(5.7), Inches(0.22), label, sz=10, bold=True, color=RED)
        y2 += Inches(0.56)

# =============================================================================
# SLIDE 18 – Key Findings
# =============================================================================
def slide_findings(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Key Findings", "Analysis")
    footer(sl, "Analysis")

    findings = [
        ("1. Stance is a fundamentally harder task than sentiment.",
         "Every model loses ~30 pp macro-F1 going from FPB to FOMC. Not about data size — about label pragmatics."),
        ("2. Sentiment is saturated by domain pretraining alone.",
         "FinBERT zero-shot reaches F1 0.9650; fine-tuning adds almost nothing on FPB. Extra capacity is spent on stance."),
        ("3. A careful fine-tuning recipe rivals domain pretraining.",
         "BERT-base with LLRD + gradual unfreeze + label smoothing matches FinBERT on sentiment and beats it on stance."),
        ("4. Multi-task training is a free lunch.",
         "One shared encoder + two heads improves single-task FinBERT by +2.07 pp sentiment F1 and +1.90 pp stance F1, and ties BERT LLRD."),
        ("5. Zero-shot sentiment ≠ stance.",
         "The native FinBERT head maps poorly to hawkish/dovish (F1 0.4874) — tone and policy direction are different signals."),
    ]
    y = BODY_Y + Inches(0.12)
    for title, body in findings:
        rect(sl, Inches(0.3), y + Inches(0.09), Inches(0.1), Inches(0.30), YEL)
        tb = sl.shapes.add_textbox(Inches(0.58), y, Inches(12.3), Inches(0.55))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = title + "  "; r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = NAV
        r2 = p.add_run(); r2.text = body;           r2.font.size = Pt(12); r2.font.color.rgb = BLACK
        y += Inches(0.62)

# =============================================================================
# SLIDE 19 – Domain Pretraining Gap chart
# =============================================================================
def slide_domain_gap(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Domain-Pretraining Gap (k=16 probe)", "Analysis")
    footer(sl, "Analysis")

    rect(sl, Inches(1.5), BODY_Y + Inches(0.1), Inches(10.0), Inches(5.0), LGREY, line=True)
    txb(sl, Inches(1.6), BODY_Y + Inches(0.3), Inches(9.8), Inches(0.4),
        "[ analysis/domain_pretraining_gap.png ]", sz=14, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)

    # mini bar chart data representation
    data = [
        ("FinBERT",     0.967, 0.487),
        ("BERT-base",   0.650, 0.374),
        ("RoBERTa",     0.672, 0.360),
    ]
    chart_x = Inches(1.8)
    chart_y = BODY_Y + Inches(0.9)
    bar_h   = Inches(0.38)
    scale   = Inches(8.5)  # width for value=1.0

    for gi, (name, sent_f1, stance_f1) in enumerate(data):
        cy = chart_y + gi * Inches(1.3)
        txb(sl, chart_x, cy, Inches(1.5), bar_h, name, sz=12, bold=True, color=NAV)
        # sentiment bar
        bw = sent_f1 * scale
        rect(sl, chart_x + Inches(1.6), cy, bw, bar_h * 0.42, NAV)
        txb(sl, chart_x + Inches(1.6) + bw + Inches(0.05), cy, Inches(1.0), bar_h * 0.42,
            f"Sent {sent_f1:.3f}", sz=10, color=NAV)
        # stance bar
        bw2 = stance_f1 * scale
        rect(sl, chart_x + Inches(1.6), cy + bar_h * 0.52, bw2, bar_h * 0.42,
             RGBColor(0xAA, 0xCC, 0xFF))
        txb(sl, chart_x + Inches(1.6) + bw2 + Inches(0.05), cy + bar_h * 0.52, Inches(1.0), bar_h * 0.42,
            f"Stance {stance_f1:.3f}", sz=10, color=NAV)

    txb(sl, Inches(0.3), FOOT_Y - Inches(0.45), Inches(12.5), Inches(0.40),
        "FinBERT outperforms BERT-base and RoBERTa by ~30 pp sentiment F1 in the low-data regime. "
        "Fine-tuning closes part of this gap — only when paired with a careful recipe (LLRD + gradual unfreeze).",
        sz=11, italic=True, color=NAV)

# =============================================================================
# SLIDE 20 – Performance Progression chart
# =============================================================================
def slide_progression(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Performance Progression Across Methods", "Analysis")
    footer(sl, "Analysis")

    rect(sl, Inches(1.5), BODY_Y + Inches(0.1), Inches(10.0), Inches(5.0), LGREY, line=True)
    txb(sl, Inches(1.6), BODY_Y + Inches(0.3), Inches(9.8), Inches(0.4),
        "[ analysis/performance_progression.png ]", sz=14, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)

    # Simple data visualisation
    models = [
        "LM Lexicon", "TF-IDF+LR", "TF-IDF+SVM", "FinBERT ZS", "FinBERT FS",
        "FinBERT FT", "BERT LLRD", "MT-FinBERT"
    ]
    sent   = [0.532, 0.823, 0.853, 0.965, 0.967, 0.946, 0.967, 0.967]
    stance = [0.389, 0.587, 0.606, 0.487, 0.453, 0.619, 0.638, 0.638]

    cx = Inches(1.8); cw = Inches(9.5); bottom = BODY_Y + Inches(5.5); scale = Inches(4.0)
    step = cw / (len(models) - 1)

    # draw lines
    for task, vals, color in [("Sentiment", sent, NAV), ("Stance", stance, RGBColor(0xAA, 0x44, 0x00))]:
        prev = None
        for i, v in enumerate(vals):
            x = cx + i * step
            y = bottom - v * scale
            r = rect(sl, x - Inches(0.06), y - Inches(0.06), Inches(0.12), Inches(0.12), color)
            if prev:
                # approximate line with narrow rect
                px, py = prev
                import math
                dx = x - px; dy = y - py
                length = math.sqrt(dx*dx + dy*dy)
                if length > 0:
                    aspect = abs(dy / dx) if dx != 0 else 99
                    mid_x = (px + x) / 2 - Inches(0.02)
                    mid_y = (py + y) / 2 - Inches(0.02)
                    rect(sl, mid_x, mid_y, Inches(0.04), Inches(0.04), color)
            prev = (x, y)
        # labels
        txb(sl, cx + (len(models)-1)*step + Inches(0.1), bottom - vals[-1]*scale - Inches(0.1),
            Inches(1.2), Inches(0.28), task, sz=9, bold=True, color=color)

    for i, m in enumerate(models):
        txb(sl, cx + i * step - Inches(0.4), bottom + Inches(0.05), Inches(0.8), Inches(0.55),
            m, sz=8, color=NAV, align=PP_ALIGN.CENTER)

    txb(sl, Inches(0.3), FOOT_Y - Inches(0.45), Inches(12.5), Inches(0.40),
        "Sentiment plateaus early (FinBERT zero-shot near-ceiling); stance keeps climbing through fine-tuning and multi-task.",
        sz=11, italic=True, color=NAV)

# =============================================================================
# SLIDE 21 – Deployment
# =============================================================================
def slide_deployment(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Deployment — CLI, Demo, and Hub Release", "Deployment")
    footer(sl, "Deployment")

    # left column
    tb = sl.shapes.add_textbox(Inches(0.4), BODY_Y + Inches(0.1), Inches(6.5), BODY_H - Inches(0.1))
    tf = tb.text_frame; tf.word_wrap = True
    para(tf, "Command-line interface", sz=14, bold=True, color=NAV)
    para(tf, '  python cli.py --text "The Committee decided to raise the target range."', sz=11, color=RGBColor(0x33, 0x33, 0x33))
    para(tf, "  Returns predicted stance and sentiment with class probabilities.", sz=12)
    para(tf, "  Loads multi-task FinBERT by default; single-task via --model flag.", sz=12)
    para(tf, "", sz=8)
    para(tf, "Interactive Gradio demo", sz=14, bold=True, color=NAV)
    para(tf, "  http://localhost:7860", sz=12, color=RGBColor(0x00, 0x55, 0xAA))
    para(tf, "  Text box + live predictions + per-class probability bars.", sz=12)
    para(tf, "  Same checkpoint as CLI — fully reproducible.", sz=12)

    # right column
    rect(sl, Inches(7.2), BODY_Y + Inches(0.1), Inches(5.7), BODY_H - Inches(0.1), LGREY, line=True)
    txb(sl, Inches(7.4), BODY_Y + Inches(0.22), Inches(5.3), Inches(0.38),
        "Hugging Face Hub — 5 public model repos", sz=14, bold=True, color=NAV)
    hub_items = [
        "Louisnguyen/finbert-fomc-stance",
        "Louisnguyen/finbert-fpb-sentiment",
        "Louisnguyen/bert-llrd-fomc-stance",
        "Louisnguyen/bert-llrd-fpb-sentiment",
        "Louisnguyen/finbert-multitask-fomc-fpb",
    ]
    hy = BODY_Y + Inches(0.72)
    for item in hub_items:
        rect(sl, Inches(7.3), hy, Inches(5.5), Inches(0.0), NAV)
        txb(sl, Inches(7.4), hy, Inches(5.4), Inches(0.36), "  " + item, sz=11, color=NAV)
        hy += Inches(0.38)
    txb(sl, Inches(7.4), hy + Inches(0.1), Inches(5.3), Inches(0.5),
        "Each repo ships weights, tokenizer config, and an inference code snippet.",
        sz=11, italic=True, color=RGBColor(0x44, 0x44, 0x44))

# =============================================================================
# SLIDE 22 – Summary
# =============================================================================
def slide_summary(prs):
    sl = prs.slides.add_slide(blank(prs))
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    title_bar(sl, "Summary", "Summary")
    footer(sl, "Summary")

    points = [
        ("Problem",
         "Joint stance (FOMC) and sentiment (FPB) classification on small, imbalanced financial corpora."),
        ("Approach",
         "Six-stage progression: LM-lexicon rules → TF-IDF → zero/few-shot probes → single-task fine-tuning (FinBERT & BERT-LLRD) → multi-task FinBERT."),
        ("Headline result",
         "Multi-task FinBERT and BERT-base LLRD tie for best on both tasks (sentiment F1 ≈ 0.967; stance F1 ≈ 0.638). Multi-task beats single-task FinBERT by +2.07 pp sentiment / +1.90 pp stance F1."),
        ("Key lesson",
         "Domain pretraining and careful fine-tuning recipes are substitutes, not complements — either one reaches the ceiling for this data size."),
    ]
    y = BODY_Y + Inches(0.25)
    for title, body in points:
        rect(sl, Inches(0.3), y, Inches(2.2), Inches(0.9), NAV)
        txb(sl, Inches(0.35), y + Inches(0.22), Inches(2.1), Inches(0.5),
            title, sz=13, bold=True, color=YEL, align=PP_ALIGN.CENTER)
        rect(sl, Inches(2.5), y, Inches(10.5), Inches(0.9), LGREY, line=True)
        txb(sl, Inches(2.6), y + Inches(0.1), Inches(10.2), Inches(0.75),
            body, sz=13, color=NAV)
        y += Inches(1.05)

    txb(sl, Inches(0.3), y + Inches(0.3), Inches(12.5), Inches(0.5),
        "Thank you — questions welcome",
        sz=22, bold=True, color=NAV, align=PP_ALIGN.CENTER)

# =============================================================================
# BUILD
# =============================================================================
def main():
    prs = new_prs()

    slide_title(prs)
    slide_outline(prs)

    section_slide(prs, "Problem")
    slide_two_tasks(prs)
    slide_why_matters(prs)

    section_slide(prs, "Datasets")
    slide_datasets(prs)
    slide_preprocessing(prs)

    section_slide(prs, "Methodology")
    slide_pipeline(prs)
    slide_baselines(prs)
    slide_pretrained(prs)
    slide_finetuning(prs)
    slide_multitask(prs)

    section_slide(prs, "Results")
    slide_baseline_results(prs)
    slide_pretrained_results(prs)
    slide_finetuned_results(prs)
    slide_multitask_result(prs)
    slide_full_results(prs)
    slide_errors(prs)

    section_slide(prs, "Analysis")
    slide_findings(prs)
    slide_domain_gap(prs)
    slide_progression(prs)

    section_slide(prs, "Deployment")
    slide_deployment(prs)

    section_slide(prs, "Summary")
    slide_summary(prs)

    out = r"D:\Study\COMP UNSW\COMP6713\financial-nlp-stance-sentiment\presentation\presentation.pptx"
    prs.save(out)
    print(f"Saved {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
