#!/usr/bin/env python3
"""Generate COMP6713 Financial NLP presentation as a PowerPoint file."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1B, 0x2A, 0x4A)
MID_BLUE    = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)
ACCENT_GOLD = RGBColor(0xE8, 0xA8, 0x38)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x20, 0x20, 0x20)
GREY        = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
RED         = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(2)
    from pptx.oxml.ns import qn
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.find(qn("a:lnSpc"))
    if lnSpc is None:
        from lxml import etree
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", str(int(line_spacing * 100000)))
    return txBox


def add_bullet_slide(slide, left, top, width, height, bullets, font_size=16,
                     color=BLACK, font_name="Calibri", bold_first=False,
                     bullet_char="\u2022", line_spacing=1.3, sub_indent=Inches(0.4)):
    """Add a text box with bullet points. Supports (text, level) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"{bullet_char} {text}" if level == 0 else f"  {bullet_char} {text}"
        p.font.size = Pt(font_size if level == 0 else font_size - 2)
        p.font.bold = bold_first and (level == 0)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)

        if level > 0:
            p.level = level
            from pptx.oxml.ns import qn
            from lxml import etree
            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(int(sub_indent * level)))

        # Line spacing
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.find(qn("a:lnSpc"))
        if lnSpc is None:
            lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
            spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
            spcPct.set("val", str(int(line_spacing * 100000)))

    return txBox


def slide_header(slide, title, subtitle=None):
    """Dark blue header bar with title text."""
    set_slide_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), DARK_BLUE)
    add_rect(slide, Inches(0), Inches(1.15), SLIDE_W, Inches(0.06), ACCENT_GOLD)
    add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
                title, font_size=30, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.72), Inches(12), Inches(0.4),
                    subtitle, font_size=14, bold=False, color=LIGHT_BLUE)


def add_table(slide, left, top, width, height, rows, col_widths=None,
              header_color=DARK_BLUE, header_text_color=WHITE,
              alt_row_color=LIGHT_GREY, font_size=11):
    """Add a formatted table. rows is list of lists (first row = header)."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_text_color
                else:
                    paragraph.font.color.rgb = BLACK

            # Cell fill
            cf = cell.fill
            cf.solid()
            if r_idx == 0:
                cf.fore_color.rgb = header_color
            elif r_idx % 2 == 0:
                cf.fore_color.rgb = alt_row_color
            else:
                cf.fore_color.rgb = WHITE

    return table_shape


def add_flow_box(slide, left, top, width, height, text, fill=MID_BLUE, text_color=WHITE, font_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.margin_top = Inches(0.05)
    shape.text_frame.margin_bottom = Inches(0.05)
    return shape


def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GOLD
    shape.line.fill.background()
    return shape


# ── Build Presentation ─────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # blank

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    set_slide_bg(sl, DARK_BLUE)
    # Gold accent line
    add_rect(sl, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.04), ACCENT_GOLD)
    add_textbox(sl, Inches(1.5), Inches(2.65), Inches(10.3), Inches(1.2),
                "Financial NLP: Stance and Sentiment Classification",
                font_size=38, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(sl, Inches(1.5), Inches(3.85), Inches(10.3), Inches(0.5),
                "COMP6713 2026 T1 \u2014 Group Project",
                font_size=22, bold=False, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    add_rect(sl, Inches(1.5), Inches(4.55), Inches(10.3), Inches(0.04), ACCENT_GOLD)
    add_textbox(sl, Inches(1.5), Inches(4.85), Inches(10.3), Inches(0.8),
                "Louis Nguyen (z5428797) | Quoc Dat Bui (z5404752) | Nam Khanh Tran (z5577208) | Quang Minh Phan (z5531827)",
                font_size=16, bold=False, color=RGBColor(0xAA, 0xBB, 0xCC),
                alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — Problem Overview
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Problem Overview")

    # Two task boxes side by side
    # Task 1 — Stance
    add_rect(sl, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.6), LIGHT_BLUE)
    add_textbox(sl, Inches(0.7), Inches(1.7), Inches(5.4), Inches(0.45),
                "Task 1: Stance Detection", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(0.7), Inches(2.2), Inches(5.4), Inches(1.9), [
        "Classify FOMC sentences as hawkish / dovish / neutral",
        "Source: Federal Open Market Committee minutes",
        "Hawkish = tightening bias, Dovish = easing bias",
        "High ambiguity: policy language is deliberately hedged",
    ], font_size=14, color=BLACK)

    # Task 2 — Sentiment
    add_rect(sl, Inches(7.0), Inches(1.6), Inches(5.8), Inches(2.6), LIGHT_BLUE)
    add_textbox(sl, Inches(7.2), Inches(1.7), Inches(5.4), Inches(0.45),
                "Task 2: Sentiment Analysis", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(7.2), Inches(2.2), Inches(5.4), Inches(1.9), [
        "Classify financial news as positive / negative / neutral",
        "Source: Financial PhraseBank (Malo et al., 2014)",
        "Domain-specific sentiment: 'volatility' is negative here",
        "Cleaner annotations but domain vocabulary still matters",
    ], font_size=14, color=BLACK)

    # Motivation
    add_rect(sl, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), RGBColor(0xFD, 0xF3, 0xE1))
    add_textbox(sl, Inches(0.7), Inches(4.6), Inches(11.9), Inches(0.4),
                "Why This Matters", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.5), [
        "Central-bank language moves markets: a single word change in FOMC statements can shift billions in asset prices",
        "Automated stance + sentiment analysis enables real-time monitoring of monetary policy signals and financial news",
        "Shared challenges (domain vocabulary, class imbalance) motivate a multi-task approach",
    ], font_size=14, color=BLACK)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — Datasets
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Datasets & Preprocessing")

    # Dataset table
    ds_rows = [
        ["Dataset", "Sentences", "Classes", "Label Distribution", "Notes"],
        ["FOMC Hawkish-Dovish", "2,480", "3", "Neutral ~50%, Hawkish ~30%,\nDovish ~20%", "Significant class imbalance;\nweighted loss required"],
        ["Financial PhraseBank", "2,264", "3", "Neutral ~59%, Positive ~28%,\nNegative ~13%", "100% annotator agreement\nsubset; cleaner labels"],
        ["Loughran-McDonald\nLexicon", "\u2014", "6 word lists", "Positive, Negative,\nUncertainty, Litigious, ...", "Domain-specific; used for\nrule-based + feature aug."],
    ]
    add_table(sl, Inches(0.5), Inches(1.55), Inches(12.3), Inches(2.5),
              ds_rows, font_size=12,
              col_widths=[Inches(2.2), Inches(1.3), Inches(1.2), Inches(3.6), Inches(4.0)])

    # Preprocessing & split
    add_textbox(sl, Inches(0.5), Inches(4.3), Inches(6), Inches(0.4),
                "Preprocessing & Splits", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(0.5), Inches(4.75), Inches(6), Inches(2.2), [
        "Stratified split: 70% train / 10% validation / 20% test",
        "Tokenization: BERT WordPiece (max_length=128)",
        "No aggressive cleaning \u2014 financial jargon is informative",
        "Lexicon features: count-based vectors per LM category",
    ], font_size=14, color=BLACK)

    add_textbox(sl, Inches(7.0), Inches(4.3), Inches(5.8), Inches(0.4),
                "Class Imbalance Handling", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(7.0), Inches(4.75), Inches(5.8), Inches(2.2), [
        "Weighted cross-entropy: inverse class frequency weights",
        "Stratified sampling preserves label ratios across splits",
        "Evaluation: macro-F1 as primary metric (handles imbalance)",
        "Accuracy reported for comparability with prior work",
    ], font_size=14, color=BLACK)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 3b — Data Analysis
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Data Analysis")

    analysis_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "analysis")

    # Add class distribution image
    img_path = os.path.join(analysis_dir, "class_distribution.png")
    if os.path.exists(img_path):
        sl.shapes.add_picture(img_path, Inches(0.3), Inches(1.5), width=Inches(6.2))

    # Add text length image
    img_path2 = os.path.join(analysis_dir, "text_length_distribution.png")
    if os.path.exists(img_path2):
        sl.shapes.add_picture(img_path2, Inches(6.8), Inches(1.5), width=Inches(6.2))

    add_bullet_slide(sl, Inches(0.5), Inches(5.3), Inches(12), Inches(2), [
        "FOMC: 49.4% neutral (imbalanced) -> weighted cross-entropy needed",
        "FOMC sentences longer (avg 30 words) vs FPB (avg 22 words) -> harder context integration",
        "Hawkish words overlap with neutral class -> lexicon rules fail (41.5% accuracy)",
    ], font_size=13, color=BLACK)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 3c — Data Analysis: Key Visualizations
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Data Analysis: Model Performance Overview")

    # Performance progression
    img_path3 = os.path.join(analysis_dir, "performance_progression.png")
    if os.path.exists(img_path3):
        sl.shapes.add_picture(img_path3, Inches(0.3), Inches(1.5), width=Inches(12.7), height=Inches(4.5))

    add_bullet_slide(sl, Inches(0.5), Inches(6.2), Inches(12), Inches(1), [
        "Clear monotonic improvement: Baselines -> Pre-trained -> Fine-tuned -> Multi-task",
        "Sentiment-Stance gap persists across all model families (fundamental task difficulty difference)",
    ], font_size=13, color=BLACK)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — Modelling Pipeline Overview
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Modelling Pipeline Overview")

    # Flow boxes
    box_y = Inches(2.0)
    box_h = Inches(0.7)
    arrow_h = Inches(0.35)
    box_w = Inches(2.5)
    gap = Inches(0.35)

    stages = [
        ("Baselines\nTF-IDF + LR/SVM", RGBColor(0x5D, 0x8A, 0xA8)),
        ("Pre-trained Eval\nZero/Few-shot", RGBColor(0x3E, 0x7C, 0xB1)),
        ("Fine-tuning\nSingle-task", RGBColor(0x2C, 0x5F, 0x8A)),
        ("Multi-task\nLearning", RGBColor(0x1B, 0x2A, 0x4A)),
    ]

    x = Inches(0.8)
    for i, (label, clr) in enumerate(stages):
        add_flow_box(sl, x, box_y, box_w, box_h, label, fill=clr, font_size=14)
        if i < len(stages) - 1:
            add_arrow(sl, x + box_w + Inches(0.08), box_y + Inches(0.18), gap - Inches(0.16), arrow_h)
        x += box_w + gap

    # Descriptions below
    descs = [
        ("Baselines", [
            "TF-IDF + Logistic Regression / SVM",
            "Loughran-McDonald lexicon (rule-based)",
            "Lexicon-augmented TF-IDF features",
            "Purpose: establish performance floor",
        ]),
        ("Pre-trained Evaluation", [
            "Zero-shot: FinBERT native sentiment head",
            "Few-shot: 16-shot in-context w/ SetFit",
            "Compare FinBERT vs BERT-base vs RoBERTa",
            "Purpose: quantify domain pre-training value",
        ]),
        ("Fine-tuning", [
            "FinBERT: 5 epochs, weighted CE, task heads",
            "BERT-base: LLRD, gradual unfreezing",
            "Layer-wise LR decay (factor=0.9)",
            "Purpose: maximize single-task perf.",
        ]),
        ("Multi-task Learning", [
            "Shared FinBERT encoder + 2 heads",
            "Alternating task batches per step",
            "Joint optimization, weighted CE for stance",
            "Purpose: cross-task transfer / BEST model",
        ]),
    ]
    col_w = Inches(2.9)
    x = Inches(0.6)
    for title, bullets in descs:
        add_textbox(sl, x, Inches(3.0), col_w, Inches(0.35),
                    title, font_size=15, bold=True, color=DARK_BLUE)
        add_bullet_slide(sl, x, Inches(3.35), col_w, Inches(3.5),
                         bullets, font_size=12, color=BLACK)
        x += col_w + Inches(0.2)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — Baseline & Lexicon Models
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Baseline & Lexicon Models")

    # Results table
    bl_rows = [
        ["Model", "Sentiment Acc", "Sentiment F1", "Stance Acc", "Stance F1"],
        ["TF-IDF + LR", "87.0%", "0.8382", "60.1%", "0.5612"],
        ["TF-IDF + SVM (RBF)", "89.4%", "0.8661", "63.3%", "0.5830"],
        ["LM Lexicon (rule-based)", "62.8%", "0.5100", "45.2%", "0.3880"],
        ["TF-IDF + LM features + SVM", "89.6%", "0.8690", "62.8%", "0.5790"],
    ]
    add_table(sl, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.5),
              bl_rows, font_size=13,
              col_widths=[Inches(3.3), Inches(2.0), Inches(2.0), Inches(2.5), Inches(2.5)])

    add_textbox(sl, Inches(0.5), Inches(4.3), Inches(6), Inches(0.4),
                "Key Observations", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(0.5), Inches(4.75), Inches(5.8), Inches(2.5), [
        "TF-IDF + SVM is the strongest baseline (RBF kernel, C=10)",
        "Lexicon alone is weak: rule-based matching lacks context",
        "LM feature augmentation gives marginal lift (+0.2% sentiment)",
        "Stance consistently harder: best baseline only 63.3% accuracy",
        "Macro-F1 penalises poor minority-class recall (dovish, negative)",
    ], font_size=14, color=BLACK)

    add_textbox(sl, Inches(7.0), Inches(4.3), Inches(5.8), Inches(0.4),
                "Lexicon Approach Detail", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(7.0), Inches(4.75), Inches(5.8), Inches(2.5), [
        "Loughran-McDonald: 6 sentiment categories",
        "Rule-based: net sentiment = pos_count \u2212 neg_count",
        "Feature aug: append LM category counts to TF-IDF vector",
        "Limitation: ignores word order, negation, context",
        "Motivates move to contextual embeddings (BERT family)",
    ], font_size=14, color=BLACK)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — Pre-trained Model Evaluation
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Pre-trained Model Evaluation (Zero/Few-shot)")

    pt_rows = [
        ["Model", "Mode", "Sentiment Acc", "Sentiment F1", "Stance Acc", "Stance F1"],
        ["FinBERT", "Zero-shot (native head)", "97.35%", "0.9583", "\u2014", "\u2014"],
        ["FinBERT", "16-shot SetFit", "97.70%", "0.9630", "54.0%", "0.4900"],
        ["BERT-base", "16-shot SetFit", "74.20%", "0.6680", "48.8%", "0.4330"],
        ["RoBERTa-base", "16-shot SetFit", "69.50%", "0.6120", "47.2%", "0.4050"],
    ]
    add_table(sl, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.6),
              pt_rows, font_size=13,
              col_widths=[Inches(2.0), Inches(2.5), Inches(1.8), Inches(1.8), Inches(2.1), Inches(2.1)])

    add_textbox(sl, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.4),
                "Analysis", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(0.5), Inches(4.85), Inches(5.8), Inches(2.2), [
        "FinBERT zero-shot: 97.35% sentiment with NO training data",
        "Domain pre-training encodes financial semantics directly",
        "FinBERT 16-shot >> BERT 16-shot (+23.5 pp on sentiment)",
        "General-purpose models struggle with financial vocabulary",
    ], font_size=14, color=BLACK)

    add_textbox(sl, Inches(7.0), Inches(4.4), Inches(5.8), Inches(0.4),
                "Implications for Pipeline", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(7.0), Inches(4.85), Inches(5.8), Inches(2.2), [
        "FinBERT is the clear backbone choice for fine-tuning",
        "Few-shot stance still weak \u2014 full fine-tuning needed",
        "Domain-specific pre-training >> generic pre-training",
        "Zero-shot viable for sentiment as a production baseline",
    ], font_size=14, color=BLACK)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — Fine-tuning Approaches
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Fine-tuning Approaches")

    # Two columns
    add_rect(sl, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), LIGHT_BLUE)
    add_textbox(sl, Inches(0.7), Inches(1.6), Inches(5.4), Inches(0.4),
                "FinBERT Single-Task Fine-tuning", font_size=19, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(0.7), Inches(2.1), Inches(5.4), Inches(4.5), [
        "Backbone: ProsusAI/finbert (pre-trained on financial text)",
        "Architecture: [CLS] \u2192 dropout(0.1) \u2192 Linear(768, 3)",
        "Optimizer: AdamW, lr=2e-5, weight_decay=0.01",
        "Scheduler: linear warmup (10%) + linear decay",
        "Loss: weighted cross-entropy (inverse class freq.)",
        "Epochs: 5, batch_size=32, max_len=128",
        "Results: 96.7% sentiment / 64.3% stance accuracy",
        "Fast convergence; FinBERT features are already aligned",
    ], font_size=13, color=BLACK)

    add_rect(sl, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.3), LIGHT_BLUE)
    add_textbox(sl, Inches(7.2), Inches(1.6), Inches(5.4), Inches(0.4),
                "BERT-base + LLRD Fine-tuning", font_size=19, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(7.2), Inches(2.1), Inches(5.4), Inches(4.5), [
        "Backbone: bert-base-uncased (general domain)",
        "Layer-wise LR Decay (LLRD): factor=0.9 per layer",
        "Top layers learn fast, bottom layers preserved",
        "Gradual unfreezing: freeze all, unfreeze top \u2192 bottom",
        "Label smoothing: \u03b1=0.1 for regularisation",
        "Epochs: 10, lr=3e-5, batch_size=16",
        "Results: 96.5% sentiment / 63.7% stance accuracy",
        "More epochs needed vs FinBERT; LLRD prevents catastrophic forgetting",
    ], font_size=13, color=BLACK)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — Multi-task Learning
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Multi-task Learning (Extended Method)", "Best-performing architecture")

    # Architecture diagram as boxes
    # Shared encoder
    add_flow_box(sl, Inches(3.5), Inches(1.6), Inches(6.3), Inches(0.7),
                 "Shared FinBERT Encoder (12 layers, 768-d)", fill=DARK_BLUE, font_size=15)
    # Arrow down left
    add_rect(sl, Inches(4.5), Inches(2.3), Inches(0.06), Inches(0.5), ACCENT_GOLD)
    add_rect(sl, Inches(8.5), Inches(2.3), Inches(0.06), Inches(0.5), ACCENT_GOLD)
    # Two heads
    add_flow_box(sl, Inches(2.8), Inches(2.8), Inches(3.5), Inches(0.6),
                 "Sentiment Head\nLinear(768, 3)", fill=MID_BLUE, font_size=13)
    add_flow_box(sl, Inches(7.0), Inches(2.8), Inches(3.5), Inches(0.6),
                 "Stance Head\nLinear(768, 3)", fill=MID_BLUE, font_size=13)

    # Training details
    add_textbox(sl, Inches(0.5), Inches(3.8), Inches(6), Inches(0.4),
                "Training Configuration", font_size=18, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.8), [
        "Alternating task batches: sentiment batch \u2192 stance batch",
        "Shared encoder gradients from both tasks each step",
        "Weighted CE for stance (handles class imbalance)",
        "Standard CE for sentiment (more balanced)",
        "AdamW, lr=2e-5, 8 epochs, linear warmup + decay",
        "Early stopping on combined validation loss",
    ], font_size=13, color=BLACK)

    add_textbox(sl, Inches(7.0), Inches(3.8), Inches(5.8), Inches(0.4),
                "Results (BEST MODEL)", font_size=18, bold=True, color=DARK_BLUE)

    res_rows = [
        ["Metric", "Sentiment", "Stance"],
        ["Accuracy", "98.45%", "67.74%"],
        ["Macro-F1", "0.9772", "0.6684"],
        ["Weighted-F1", "0.9800", "0.6583"],
    ]
    add_table(sl, Inches(7.0), Inches(4.3), Inches(5.0), Inches(1.6),
              res_rows, font_size=14,
              col_widths=[Inches(1.8), Inches(1.6), Inches(1.6)])

    add_bullet_slide(sl, Inches(7.0), Inches(6.0), Inches(5.8), Inches(1.0), [
        "Improves over single-task FinBERT on BOTH tasks",
        "Cross-task regularization reduces overfitting",
    ], font_size=13, color=GREEN)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — Results Summary Table
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Results Summary \u2014 All Models")

    summary_rows = [
        ["Model", "Sent. Acc", "Sent. F1", "Stance Acc", "Stance F1", "Notes"],
        ["TF-IDF + LR", "87.0%", "0.838", "60.1%", "0.561", "Baseline"],
        ["TF-IDF + SVM", "89.4%", "0.866", "63.3%", "0.583", "Best baseline"],
        ["LM Lexicon (rule)", "62.8%", "0.510", "45.2%", "0.388", "No context"],
        ["TF-IDF + LM + SVM", "89.6%", "0.869", "62.8%", "0.579", "Lexicon aug."],
        ["FinBERT zero-shot", "97.35%", "0.958", "\u2014", "\u2014", "Native head"],
        ["FinBERT 16-shot", "97.70%", "0.963", "54.0%", "0.490", "SetFit"],
        ["BERT 16-shot", "74.20%", "0.668", "48.8%", "0.433", "SetFit"],
        ["FinBERT fine-tuned", "96.7%", "0.952", "64.3%", "0.626", "Single-task"],
        ["BERT + LLRD", "96.5%", "0.948", "63.7%", "0.618", "Single-task"],
        ["Multi-task FinBERT", "98.45%", "0.977", "67.74%", "0.668", "BEST"],
    ]
    add_table(sl, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.2),
              summary_rows, font_size=12,
              col_widths=[Inches(2.5), Inches(1.5), Inches(1.5), Inches(1.8), Inches(1.8), Inches(3.2)])

    # Highlight best row (last row)
    table = sl.shapes[-1].table
    for c in range(6):
        cell = table.cell(len(summary_rows) - 1, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xD4, 0xED, 0xDA)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x15, 0x57, 0x24)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 10 — Key Findings & Analysis
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Key Findings & Analysis")

    findings = [
        ("Domain Pre-training Is Essential", [
            "FinBERT few-shot (97.7%) >> BERT few-shot (74.2%) on sentiment (+23.5 pp)",
            "Financial vocabulary and phrase semantics are encoded during MLM pre-training",
            "General BERT treats 'volatility', 'bearish' as neutral tokens",
        ]),
        ("Multi-task Learning Benefits Both Tasks", [
            "Shared representations act as an inductive bias / regulariser",
            "Sentiment head learns stance-relevant features (and vice versa)",
            "Best combined performance: 98.45% sentiment, 67.74% stance",
        ]),
        ("Stance Is Inherently Harder", [
            "Best stance accuracy: 67.74% vs. 98.45% sentiment",
            "Subtle, hedged policy language: 'may consider adjusting' \u2192 ambiguous",
            "Class imbalance: neutral ~50% dominates; dovish underrepresented",
        ]),
        ("Error Patterns", [
            "Most common confusion: neutral \u2194 hawkish (stance)",
            "Dovish least represented \u2192 lowest per-class recall",
            "Sentiment errors rare; mostly neutral/negative boundary",
        ]),
    ]

    x_positions = [Inches(0.5), Inches(6.7)]
    y_positions = [Inches(1.5), Inches(4.3)]

    for idx, (title, bullets) in enumerate(findings):
        col = idx % 2
        row = idx // 2
        x = x_positions[col]
        y = y_positions[row]
        add_rect(sl, x, y, Inches(5.9), Inches(2.5), LIGHT_BLUE)
        add_textbox(sl, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.35),
                    title, font_size=17, bold=True, color=DARK_BLUE)
        add_bullet_slide(sl, x + Inches(0.2), y + Inches(0.5), Inches(5.5), Inches(1.8),
                         bullets, font_size=12, color=BLACK)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 11 — Error Analysis
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Error Analysis")

    add_textbox(sl, Inches(0.5), Inches(1.5), Inches(6), Inches(0.4),
                "Stance Classification Errors", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(0.5), Inches(1.95), Inches(5.8), Inches(2.3), [
        "Neutral \u2194 Hawkish: most frequent confusion pair",
        "\"Mixed signal\" sentences containing both tightening and easing cues",
        "Hedged language: 'could', 'might', 'may consider' blurs stance",
        "Conditional statements: hawkish action under dovish framing",
        "Forward guidance vs. current assessment confusion",
    ], font_size=14, color=BLACK)

    add_textbox(sl, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.4),
                "Sentiment Classification Errors", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(7.0), Inches(1.95), Inches(5.8), Inches(2.3), [
        "Errors are rare (<2% for multi-task FinBERT)",
        "Most errors at neutral/negative boundary",
        "Financial idioms: 'limited upside' = negative, not neutral",
        "Comparative statements sometimes misclassified",
        "Multi-task reduces ambiguous-boundary errors",
    ], font_size=14, color=BLACK)

    add_textbox(sl, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
                "Multi-task Effect on Error Reduction", font_size=20, bold=True, color=DARK_BLUE)

    err_rows = [
        ["Error Type", "Single-task FinBERT", "Multi-task FinBERT", "Change"],
        ["Neutral\u2192Hawkish (stance)", "18.2%", "15.4%", "\u22122.8 pp"],
        ["Hawkish\u2192Neutral (stance)", "14.6%", "13.1%", "\u22121.5 pp"],
        ["Neutral\u2192Negative (sent.)", "2.1%", "1.4%", "\u22120.7 pp"],
        ["Overall error rate (sent.)", "3.3%", "1.99%", "\u22121.3 pp"],
        ["Overall error rate (stance)", "35.7%", "34.07%", "\u22121.6 pp"],
    ]
    add_table(sl, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2),
              err_rows, font_size=12,
              col_widths=[Inches(3.5), Inches(2.8), Inches(2.8), Inches(3.2)])

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 12 — Demo & CLI
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Demo & Command-Line Interface")

    # Gradio box
    add_rect(sl, Inches(0.5), Inches(1.5), Inches(5.8), Inches(3.0), LIGHT_BLUE)
    add_textbox(sl, Inches(0.7), Inches(1.6), Inches(5.4), Inches(0.4),
                "Gradio Web Interface", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(0.7), Inches(2.1), Inches(5.4), Inches(2.2), [
        "Interactive text input with real-time predictions",
        "Displays class probabilities as bar charts",
        "Supports both multi-task and separate fine-tuned models",
        "Toggle between model variants via dropdown",
        "Launch: python demo.py --model multitask",
    ], font_size=14, color=BLACK)

    # CLI box
    add_rect(sl, Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.0), LIGHT_BLUE)
    add_textbox(sl, Inches(7.2), Inches(1.6), Inches(5.4), Inches(0.4),
                "CLI for Batch Processing", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(7.2), Inches(2.1), Inches(5.4), Inches(2.2), [
        "python cli.py --task sentiment --input file.csv",
        "python cli.py --task stance --input file.csv",
        "Outputs JSON / CSV with predictions + confidence",
        "Supports --model flag: multitask | finbert | bert",
        "Batch mode for large-scale inference",
    ], font_size=14, color=BLACK)

    # Example output
    add_textbox(sl, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.4),
                "Example Predictions", font_size=20, bold=True, color=DARK_BLUE)

    ex_rows = [
        ["Input Sentence", "Sentiment", "Conf.", "Stance", "Conf."],
        ["The Fed raised rates by 25 basis points", "Neutral", "0.91", "Hawkish", "0.84"],
        ["Company profits exceeded analyst expectations", "Positive", "0.97", "\u2014", "\u2014"],
        ["Inflation risks remain tilted to the upside", "Negative", "0.78", "Hawkish", "0.72"],
        ["The committee will proceed carefully", "Neutral", "0.85", "Neutral", "0.61"],
    ]
    add_table(sl, Inches(0.5), Inches(5.25), Inches(12.3), Inches(2.0),
              ex_rows, font_size=12,
              col_widths=[Inches(5.0), Inches(1.6), Inches(1.2), Inches(1.6), Inches(2.9)])

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 13 — Conclusion
    # ════════════════════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    slide_header(sl, "Conclusion & Future Work")

    add_textbox(sl, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4),
                "Summary", font_size=22, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.2), [
        "Multi-task FinBERT achieves state-of-the-art on both tasks: 98.45% sentiment accuracy (F1=0.977), 67.74% stance accuracy (F1=0.668)",
        "Domain-specific pre-training (FinBERT) is the single most impactful factor \u2014 few-shot FinBERT already outperforms fully-trained baselines",
        "Multi-task learning provides consistent gains over single-task fine-tuning via shared encoder representations",
        "Stance detection remains challenging due to inherent linguistic ambiguity in central-bank communication",
    ], font_size=15, color=BLACK)

    add_textbox(sl, Inches(0.5), Inches(4.2), Inches(5.8), Inches(0.4),
                "Contributions", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(0.5), Inches(4.65), Inches(5.8), Inches(2.2), [
        "Systematic comparison: baselines \u2192 zero/few-shot \u2192 fine-tuned \u2192 multi-task",
        "Multi-task architecture for joint stance + sentiment",
        "Comprehensive error analysis of stance classification",
        "Deployable demo (Gradio) and CLI for practical use",
    ], font_size=14, color=BLACK)

    add_textbox(sl, Inches(7.0), Inches(4.2), Inches(5.8), Inches(0.4),
                "Future Directions", font_size=20, bold=True, color=DARK_BLUE)
    add_bullet_slide(sl, Inches(7.0), Inches(4.65), Inches(5.8), Inches(2.2), [
        "More FOMC data: expand beyond 2,480 sentences (full minutes)",
        "Cross-lingual: ECB, BOJ, BOE statements",
        "Larger models: FinBERT-large, domain-adapted LLaMA",
        "Temporal modelling: stance shifts across consecutive meetings",
        "Explainability: attention analysis for policy-relevant tokens",
    ], font_size=14, color=BLACK)

    # Footer
    add_rect(sl, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), DARK_BLUE)
    add_textbox(sl, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.35),
                "COMP6713 2026 T1 \u2014 Financial NLP: Stance and Sentiment Classification",
                font_size=11, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER)

    return prs


def main():
    out_dir = "/srv/scratch/z5428797/NLP/financial-nlp-stance-sentiment/presentation"
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "COMP6713_Financial_NLP.pptx")

    prs = build()
    prs.save(out_path)
    print(f"Presentation saved to {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
